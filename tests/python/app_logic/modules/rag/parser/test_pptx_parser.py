import tempfile
from pathlib import Path
from pptx import Presentation
import pytest
from ataraxai.app_logic.modules.rag.parser.pptx_parser import PPTXParser
from ataraxai.app_logic.modules.rag.parser.document_base_parser import DocumentChunk

def create_sample_pptx(slides_content):
    prs = Presentation()
    for content in slides_content:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank slide
        for text in content:
            txBox = slide.shapes.add_textbox(left=0, top=0, width=100, height=100)
            tf = txBox.text_frame
            tf.text = text
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    tmp.close()
    return Path(tmp.name)

def test_parse_single_slide_single_text(monkeypatch):
    pptx_path = create_sample_pptx([["Hello World"]])
    parser = PPTXParser()
    chunks = parser.parse(pptx_path)
    assert len(chunks) == 1
    assert isinstance(chunks[0], DocumentChunk)
    assert chunks[0].content == "Hello World"
    assert chunks[0].source == str(pptx_path)
    assert chunks[0].metadata == {"slide": 1}
    pptx_path.unlink()

def test_parse_multiple_slides_multiple_texts():
    pptx_path = create_sample_pptx([
        ["Slide 1 Text A", "Slide 1 Text B"],
        ["Slide 2 Text"],
        [],
        ["Slide 4 Text A", "Slide 4 Text B", "Slide 4 Text C"]
    ])
    parser = PPTXParser()
    chunks = parser.parse(pptx_path)
    assert len(chunks) == 3
    assert chunks[0].content == "Slide 1 Text A\nSlide 1 Text B"
    assert chunks[0].metadata == {"slide": 1}
    assert chunks[1].content == "Slide 2 Text"
    assert chunks[1].metadata == {"slide": 2}
    assert chunks[2].content == "Slide 4 Text A\nSlide 4 Text B\nSlide 4 Text C"
    assert chunks[2].metadata == {"slide": 4}
    pptx_path.unlink()

def test_parse_empty_pptx():
    prs = Presentation()
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    tmp.close()
    pptx_path = Path(tmp.name)
    parser = PPTXParser()
    chunks = parser.parse(pptx_path)
    assert chunks == []
    pptx_path.unlink()

def test_parse_slide_with_empty_textbox():
    pptx_path = create_sample_pptx([["", "   ", "Valid Text"]])
    parser = PPTXParser()
    chunks = parser.parse(pptx_path)
    assert len(chunks) == 1
    assert chunks[0].content == "Valid Text"
    pptx_path.unlink()