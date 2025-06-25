from pptx import Presentation
from pathlib import Path
from typing import List
from ataraxai.app_logic.modules.rag.parser.document_base_parser import (
    DocumentParser,
    DocumentChunk,
)


class PPTXParser(DocumentParser):
    def parse(self, path: Path) -> List[DocumentChunk]:
        """
        Parses a PowerPoint (.pptx) file and extracts text content from each slide.

        Args:
            path (Path): The file path to the PowerPoint presentation.

        Returns:
            List[DocumentChunk]: A list of DocumentChunk objects, each containing the concatenated text from a slide,
            the source file path, and metadata indicating the slide number.
        """
        prs = Presentation(str(path))
        chunks = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                chunks.append(
                    DocumentChunk(
                        content="\n".join(texts),
                        source=str(path),
                        metadata={"slide": i + 1},
                    )
                )
        return chunks
