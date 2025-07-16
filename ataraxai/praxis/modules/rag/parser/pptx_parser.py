from pptx import Presentation
from pathlib import Path
from typing import List
from ataraxai.praxis.modules.rag.parser.document_base_parser import (
    DocumentParser,
    DocumentChunk,
)


class PPTXParser(DocumentParser):
    def parse(self, path: Path) -> List[DocumentChunk]:
        """
        Parses a PowerPoint (.pptx) file and extracts text content from each slide.

        Args:
            path (Path): The file path to the PowerPoint presentation.

        Returns:
            List[DocumentChunk]: A list of DocumentChunk objects, each containing the concatenated text from a slide,
            the source file path, and metadata indicating the slide number.
        """
        prs = Presentation(str(path))
        chunks: List[DocumentChunk] = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():  # type: ignore
                    texts.append(shape.text.strip())  # type: ignore
            if texts:
                chunks.append(
                    DocumentChunk(
                        content="\n".join(texts),
                        source=str(path),
                        metadata={"slide": i + 1},
                    )
                )
        return chunks


if __name__ == "__main__":
    pptx_path = Path(
        "tests/python/assets/Attention is All you need.pptx"
    )
    parser = PPTXParser()
    document_chunks = parser.parse(pptx_path)
    for chunk in document_chunks:
        print(f"Slide {chunk.metadata.get('slide')}: {chunk.content[:100]}...")
        print(f"Source: {chunk.source}")
        print(f"Metadata: {chunk.metadata}")
    print(f"Total slides parsed: {len(document_chunks)}")
